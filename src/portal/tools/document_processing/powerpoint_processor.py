"""
PowerPoint Processor Tool
==========================

Create and edit PowerPoint (PPTX) presentations programmatically.

Features:
- Create presentations from scratch
- Add slides with different layouts
- Insert text, images, charts, tables
- Apply themes and formatting
- Convert from Markdown/HTML
- Extract content

Install: pip install python-pptx Pillow
"""

import logging
from pathlib import Path
from typing import Any

from portal.core.interfaces.tool import BaseTool, ToolCategory, ToolMetadata, ToolParameter

logger = logging.getLogger(__name__)

# Check dependencies
try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PowerPointProcessorTool(BaseTool):
    """
    Create and edit PowerPoint presentations.

    Perfect for generating reports, pitch decks, and presentations programmatically.
    """

    # Slide layouts
    LAYOUTS = {
        "title": 0,
        "title_content": 1,
        "section_header": 2,
        "two_content": 3,
        "comparison": 4,
        "title_only": 5,
        "blank": 6,
        "content_caption": 7,
        "picture_caption": 8,
    }

    def __init__(self) -> None:
        super().__init__()

    def _get_metadata(self) -> ToolMetadata:
        return ToolMetadata(
            name="powerpoint_processor",
            description="Create and edit PowerPoint (PPTX) presentations",
            category=ToolCategory.UTILITY,
            version="1.0.0",
            requires_confirmation=False,
            parameters=[
                ToolParameter(
                    name="action",
                    param_type="string",
                    description="Action: create, add_slide, add_text, add_image, add_chart, read, save",
                    required=True,
                ),
                ToolParameter(
                    name="file_path",
                    param_type="string",
                    description="Path to PowerPoint file",
                    required=True,
                ),
                ToolParameter(
                    name="layout",
                    param_type="string",
                    description="Slide layout: title, title_content, section_header, two_content, blank",
                    required=False,
                    default="title_content",
                ),
                ToolParameter(
                    name="title", param_type="string", description="Slide title", required=False
                ),
                ToolParameter(
                    name="content",
                    param_type="string",
                    description="Slide content (text or bullet points)",
                    required=False,
                ),
                ToolParameter(
                    name="bullet_points",
                    param_type="list",
                    description="List of bullet points",
                    required=False,
                ),
                ToolParameter(
                    name="image_path",
                    param_type="string",
                    description="Path to image file",
                    required=False,
                ),
                ToolParameter(
                    name="image_position",
                    param_type="object",
                    description="Image position: {left, top, width, height} in inches",
                    required=False,
                ),
                ToolParameter(
                    name="chart_data",
                    param_type="object",
                    description="Chart data for creating charts",
                    required=False,
                ),
                ToolParameter(
                    name="chart_type",
                    param_type="string",
                    description="Chart type: bar, line, pie",
                    required=False,
                ),
                ToolParameter(
                    name="theme",
                    param_type="object",
                    description="Theme settings: {background_color, title_color, text_color}",
                    required=False,
                ),
            ],
        )

    async def execute(self, parameters: dict[str, Any]) -> dict[str, Any]:
        """Execute PowerPoint operation"""

        if not PPTX_AVAILABLE:
            return self._error_response(
                "python-pptx not installed. Install: pip install python-pptx Pillow"
            )

        action = parameters.get("action", "").lower()
        dispatch = {
            "create": self._create_presentation,
            "add_slide": self._add_slide,
            "add_image": self._add_image_to_slide,
            "add_chart": self._add_chart_to_slide,
            "read": self._read_presentation,
            "save": self._save_presentation,
        }
        handler = dispatch.get(action)
        if handler is None:
            return self._error_response(f"Unknown action: {action}")
        return await handler(parameters)

    async def _create_presentation(self, parameters: dict[str, Any]) -> dict[str, Any]:
        """Create new presentation"""

        file_path = Path(parameters.get("file_path", "")).expanduser()
        title = parameters.get("title", "New Presentation")
        subtitle = parameters.get("content", "")

        try:
            # Create presentation
            prs = Presentation()

            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)

            # Set title and subtitle
            slide.shapes.title.text = title
            if subtitle and len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle

            # Save
            file_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(file_path)

            return self._success_response(
                result={"file_path": str(file_path), "slides": 1}, metadata={"title": title}
            )

        except Exception as e:
            logger.error("PowerPoint creation error: %s", e)
            return self._error_response(f"Creation error: {e}")

    async def _add_slide(self, parameters: dict[str, Any]) -> dict[str, Any]:
        """Add slide to presentation"""

        file_path = Path(parameters.get("file_path", "")).expanduser()
        layout_name = parameters.get("layout", "title_content")
        title = parameters.get("title", "")
        content = parameters.get("content", "")
        bullet_points = parameters.get("bullet_points", [])

        if not file_path.exists():
            return self._error_response(f"File not found: {file_path}")

        try:
            # Load presentation
            prs = Presentation(file_path)

            # Get layout
            layout_idx = self.LAYOUTS.get(layout_name, 1)
            slide_layout = prs.slide_layouts[layout_idx]

            # Add slide
            slide = prs.slides.add_slide(slide_layout)

            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = title

            # Add content
            if content or bullet_points:
                # Find content placeholder
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 2:  # Body placeholder
                        text_frame = shape.text_frame
                        text_frame.clear()

                        if bullet_points:
                            for point in bullet_points:
                                p = text_frame.add_paragraph()
                                p.text = point
                                p.level = 0
                        elif content:
                            text_frame.text = content

                        break

            # Save
            prs.save(file_path)

            return self._success_response(
                result={"slides_total": len(prs.slides)},
                metadata={"slide_added": title or "Untitled"},
            )

        except Exception as e:
            logger.error("PowerPoint add slide error: %s", e)
            return self._error_response(f"Add slide error: {e}")

    async def _add_image_to_slide(self, parameters: dict[str, Any]) -> dict[str, Any]:
        """Add image to slide"""

        file_path = Path(parameters.get("file_path", "")).expanduser()
        image_path = Path(parameters.get("image_path", "")).expanduser()
        slide_idx = parameters.get("slide_index", -1)  # -1 = last slide
        position = parameters.get("image_position", {})

        if not file_path.exists():
            return self._error_response(f"File not found: {file_path}")

        if not image_path.exists():
            return self._error_response(f"Image not found: {image_path}")

        try:
            # Load presentation
            prs = Presentation(file_path)

            # Get slide
            if slide_idx == -1:
                slide = prs.slides[-1]
            else:
                if slide_idx >= len(prs.slides):
                    return self._error_response(f"Slide index {slide_idx} out of range")
                slide = prs.slides[slide_idx]

            # Add image
            left = Inches(position.get("left", 2))
            top = Inches(position.get("top", 2))
            width = Inches(position.get("width", 4))
            height = Inches(position.get("height", 3))

            slide.shapes.add_picture(str(image_path), left, top, width, height)

            # Save
            prs.save(file_path)

            return self._success_response(
                result={"image_added": str(image_path)},
                metadata={"slide_index": slide_idx if slide_idx != -1 else len(prs.slides) - 1},
            )

        except Exception as e:
            logger.error("PowerPoint add image error: %s", e)
            return self._error_response(f"Add image error: {e}")

    async def _add_chart_to_slide(self, parameters: dict[str, Any]) -> dict[str, Any]:
        """Add chart to slide"""

        file_path = Path(parameters.get("file_path", "")).expanduser()
        chart_data = parameters.get("chart_data", {})
        chart_type = parameters.get("chart_type", "bar")
        slide_idx = parameters.get("slide_index", -1)

        if not file_path.exists():
            return self._error_response(f"File not found: {file_path}")

        if not chart_data:
            return self._error_response("No chart data provided")

        try:
            # Load presentation
            prs = Presentation(file_path)

            # Get slide
            if slide_idx == -1:
                slide = prs.slides[-1]
            else:
                slide = prs.slides[slide_idx]

            # Prepare chart data
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = chart_data.get("categories", [])

            for series in chart_data.get("series", []):
                chart_data_obj.add_series(series["name"], series["values"])

            # Add chart
            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)

            if chart_type == "bar":
                chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
            elif chart_type == "line":
                chart_type_enum = XL_CHART_TYPE.LINE
            elif chart_type == "pie":
                chart_type_enum = XL_CHART_TYPE.PIE
            else:
                return self._error_response(f"Unknown chart type: {chart_type}")

            slide.shapes.add_chart(chart_type_enum, x, y, cx, cy, chart_data_obj)

            # Save
            prs.save(file_path)

            return self._success_response(
                result={"chart_added": chart_type},
                metadata={"categories": len(chart_data.get("categories", []))},
            )

        except Exception as e:
            logger.error("PowerPoint add chart error: %s", e)
            return self._error_response(f"Add chart error: {e}")

    async def _read_presentation(self, parameters: dict[str, Any]) -> dict[str, Any]:
        """Read presentation content"""

        file_path = Path(parameters.get("file_path", "")).expanduser()

        if not file_path.exists():
            return self._error_response(f"File not found: {file_path}")

        try:
            # Load presentation
            prs = Presentation(file_path)

            # Extract content
            slides_content = []

            for idx, slide in enumerate(prs.slides):
                slide_data = {"index": idx, "title": "", "content": [], "notes": ""}

                # Get title
                if slide.shapes.title:
                    slide_data["title"] = slide.shapes.title.text

                # Get text content
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape != slide.shapes.title and shape.text:
                            slide_data["content"].append(shape.text)

                # Get notes
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    slide_data["notes"] = notes_slide.notes_text_frame.text

                slides_content.append(slide_data)

            return self._success_response(
                result={"slides": slides_content},
                metadata={
                    "total_slides": len(prs.slides),
                    "slide_width": prs.slide_width,
                    "slide_height": prs.slide_height,
                },
            )

        except Exception as e:
            logger.error("PowerPoint read error: %s", e)
            return self._error_response(f"Read error: {e}")

    async def _save_presentation(self, parameters: dict[str, Any]) -> dict[str, Any]:
        """Save presentation (placeholder for future modifications)"""

        file_path = Path(parameters.get("file_path", "")).expanduser()

        if not file_path.exists():
            return self._error_response(f"File not found: {file_path}")

        return self._success_response(result={"saved": str(file_path)}, metadata={})
